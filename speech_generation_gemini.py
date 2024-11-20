from flask import Flask, request, jsonify
from flask_cors import CORS
import logging
import pyttsx3
from PyPDF2 import PdfReader
from pptx import Presentation
import google.generativeai as genai
import os

app = Flask(__name__)
CORS(app)

# Set up logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

# Configure Gemini API
genai.configure(api_key='AIzaSyBbMqcZP9FSOTB7yLCNkVXMu7XpJ5JefEs')
model = genai.GenerativeModel('gemini-pro')

def process_pdf(file):
    try:
        logger.info("Processing PDF file")
        reader = PdfReader(file)
        text = []
        
        for page in reader.pages:
            content = page.extract_text()
            if content:
                # Clean up the extracted text
                content = content.strip()
                content = ' '.join(content.split())  # Remove extra whitespace
                text.append(content)
        
        processed_text = ' '.join(text)
        logger.info(f"Successfully extracted {len(text)} pages from PDF")
        logger.debug(f"Extracted text sample: {processed_text[:500]}...")  # Debug log
        
        if not processed_text:
            raise ValueError("No text could be extracted from the PDF")
            
        return processed_text
        
    except Exception as e:
        logger.error(f"Error processing PDF: {str(e)}")
        raise

def process_pptx(file):
    try:
        logger.info("Processing PPTX file")
        prs = Presentation(file)
        text = []
        
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from all possible shape types
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                    
                # Handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_text.append(cell.text.strip())
                                
            if slide_text:
                text.append(f"Slide {slide_number}: {' '.join(slide_text)}")
        
        processed_text = ' '.join(text)
        logger.info(f"Successfully extracted {len(text)} slides from PPTX")
        logger.debug(f"Extracted text sample: {processed_text[:500]}...")  # Debug log
        
        if not processed_text:
            raise ValueError("No text could be extracted from the PPTX")
            
        return processed_text
        
    except Exception as e:
        logger.error(f"Error processing PPTX: {str(e)}")
        raise

def generate_summary(text):
    try:
        logger.info("Generating summary from text")
        logger.debug(f"Input text sample: {text[:500]}...")  # Debug log
        
        prompt = f"""
        Act as an expert teacher. Your task is to create a clear, engaging summary of the following content.
        Make it conversational, as if you're explaining it to a student.
        
        Guidelines:
        1. Start with a brief introduction of the topic
        2. Break down the main concepts
        3. Use simple, clear language
        4. Keep it engaging and natural
        5. Focus on the most important points
        6. End with a brief conclusion
        
        Content to summarize:
        {text}
        """
        
        response = model.generate_content(prompt)
        
        if not response or not response.text:
            raise ValueError("No response received from Gemini API")
            
        summary = response.text
        
        # Clean up the text for speech
        summary = summary.replace('*', '')
        summary = summary.replace('#', '')
        summary = summary.replace('\n', ' ')
        summary = ' '.join(summary.split())  # Remove extra whitespace
        
        logger.info("Successfully generated summary")
        logger.debug(f"Generated summary sample: {summary[:500]}...")  # Debug log
        
        return summary
        
    except Exception as e:
        logger.error(f"Error generating summary: {str(e)}")
        raise

@app.route('/process-file', methods=['POST'])
def process_file():
    logger.debug('Received file upload request')
    
    if 'file' not in request.files:
        logger.error('No file in request')
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        logger.error('Empty filename')
        return jsonify({'error': 'No file selected'}), 400

    try:
        logger.info(f'Processing file: {file.filename}')
        
        # Extract text based on file type
        if file.filename.lower().endswith('.pdf'):
            text = process_pdf(file)
        elif file.filename.lower().endswith('.pptx'):
            text = process_pptx(file)
        elif file.filename.lower().endswith('.txt'):
            text = file.read().decode('utf-8')
            text = text.strip()
            logger.debug(f"Extracted text from TXT: {text[:500]}...")  # Debug log
        else:
            logger.error(f'Unsupported file type: {file.filename}')
            return jsonify({'error': 'Unsupported file type. Please upload PDF, PPTX, or TXT files only.'}), 400

        if not text or len(text.strip()) == 0:
            raise ValueError("No text could be extracted from the file")

        # Generate summary using Gemini
        logger.info('Generating summary')
        summary = generate_summary(text)
        
        if not summary:
            raise ValueError("Failed to generate summary")
        
        logger.info('Successfully processed file')
        return jsonify({'summary': summary})

    except Exception as e:
        error_message = str(e)
        logger.error(f'Error processing file: {error_message}', exc_info=True)
        return jsonify({'error': f'Error processing file: {error_message}'}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5000)
